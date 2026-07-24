"""
ingestion.py — walk a KB's RAG source, chunk + embed changed files, and
upsert them into the vector store (Phase 8a). Also the text-extraction +
chunk/embed/upsert engine behind the file-upload path (Phase 10, see
`extract_text` / `ingest_text_document` below and main.py's
POST /kb-repos/{id}/files).

`ingest_kb(descriptor)` is the single entry point for repo ingestion (also
reachable via `python -m ingest --kb <id>` and per-KB `/admin/reload`,
providers/rag.py's RagProvider.reload). It:

  1. Resolves the descriptor's rag source (RagSource directly for a rag-only
     KB, or HybridSource.rag for a hybrid one). Returns a zero-count no-op if
     there isn't one — this is how customer_sap (hybrid, `source.rag: null`
     in this phase) stays completely un-ingested and byte-identical.
  2. If the source has a `git_url` (an Azure DevOps Repos URL), shallow-clones
     it into a temp dir; for a private repo the PAT is read from
     `os.environ[auth_token_env]` and passed via the Azure-recommended
     `-c http.extraHeader="Authorization: Basic <base64(':'+PAT)>"` git
     config flag — never embedded in the URL and never logged (see
     `_clone_git_repo`). The temp clone dir is always removed in a `finally`.
     Otherwise walks the local `roots` (relative to backend/, unless a root
     is already absolute — tests point `roots` straight at a tmp dir).
  3. Walks files per `include_globs`/`exclude_globs`, loads text (md/code
     read as UTF-8, `.pdf` via pypdf), and skips any file whose sha256 has not
     changed since the last successful ingest (tracked in `kb_documents`).
  4. Chunks changed files into ~800-word overlapping windows, batch-embeds
     via `embeddings.embed_texts`, and upserts into the configured
     `VectorStore` (replacing that document's previous chunks, if any) — the
     chunk/embed/upsert/replace step is `ingest_text_document`, shared with
     the upload path so both write `kb_documents`/`kb_chunks` identically.
  5. Prunes vanished files: after the walk, every `kb_documents` row for this
     kb_id whose path was NOT encountered — a file deleted from the repo, or
     excluded by a narrowed `include_globs`/`exclude_globs` — is deleted
     along with its `kb_chunks`, EXCEPT any document whose path starts with
     `uploads/` (user-uploaded files, populated by main.py's
     POST /kb-repos/{id}/files, never touched by this walk). Only runs after
     a successful walk — a failed clone raises before this point, leaving
     prior data untouched.

Returns `{"documents": <new/changed>, "chunks": <embedded+upserted>,
"skipped": <unchanged>, "removed": <pruned vanished/excluded documents>}`.
"""

from __future__ import annotations

import base64
import fnmatch
import hashlib
import io
import logging
import os
import shutil
import subprocess
import tempfile
from datetime import datetime, timezone
from pathlib import Path

from kb._schema import HybridSource, KBDescriptor, RagSource

log = logging.getLogger(__name__)

_BACKEND_DIR = Path(__file__).resolve().parent

# ~800-token chunks, word-count approximation (a simple char/word split is
# enough for this phase — no tokenizer dependency).
_CHUNK_WORDS = 800
_CHUNK_OVERLAP_WORDS = 100

_TEXT_SUFFIXES_HINT = {".md", ".txt", ".py", ".yaml", ".yml", ".json", ".rst", ".sql", ".cfg", ".ini", ".toml"}
_PDF_SUFFIX = ".pdf"

# Never descend into VCS metadata regardless of the descriptor's exclude_globs.
_ALWAYS_SKIP_DIR_NAMES = {".git", ".hg", ".svn"}


def _utcnow() -> datetime:
    return datetime.now(timezone.utc)


# ── Source resolution ───────────────────────────────────────────────────────


def _resolve_rag_source(descriptor: KBDescriptor) -> RagSource | None:
    """A rag-only KB's source *is* the RagSource; a hybrid KB's is optional
    (HybridSource.rag); a structured-only KB has none."""
    source = descriptor.source
    if isinstance(source, RagSource):
        return source
    if isinstance(source, HybridSource):
        return source.rag
    return None


def _resolve_root(root_str: str) -> Path:
    p = Path(root_str)
    return p if p.is_absolute() else (_BACKEND_DIR / p)


# ── Azure DevOps Repos clone (git subprocess, fully mockable) ──────────────


# Curated, safe reasons for a `git clone` failure aimed at a non-technical
# stakeholder — matched against exc.stderr (lowercased) keyword-by-keyword,
# first match wins. Never the raw stderr (which may echo back the failing
# URL) and never `cmd` (carries the auth header). Exposed so main.py's
# _friendly_ingest_error can recognize an already-friendly clone message and
# pass it through unchanged instead of re-wrapping it.
CLONE_ERROR_KEYWORDS: tuple[tuple[tuple[str, ...], str], ...] = (
    (
        ("authentication failed", "could not read username", "invalid username or password", "403"),
        "Couldn't access the repository — check the access token (private repos require one).",
    ),
    (
        ("repository not found", "not found", "does not exist", "404"),
        "Repository not found — check the Git URL and branch.",
    ),
    # Checked BEFORE the network bucket: git reports these as
    # "warning: unable to access '<path>': Filename too long", so the network
    # bucket must NOT key on the (far too broad) "unable to access" — git
    # prefixes it to path, TLS, and network failures alike.
    (
        ("filename too long", "cannot create directory", "unable to create file", "checkout failed"),
        "The repository downloaded, but its files couldn't be unpacked here — usually a path that's too long for this system (a deeply-nested folder). Try a repo with shorter paths.",
    ),
    (
        ("ssl certificate", "certificate verify failed", "self-signed certificate", "unable to get local issuer"),
        "Couldn't verify the repository's security certificate — usually a corporate network/proxy. Ask IT to trust the proxy's certificate for git.",
    ),
    (
        ("could not resolve host", "connection", "timed out", "network is unreachable"),
        "Couldn't reach the repository host — check the URL and your network.",
    ),
)

_CLONE_ERROR_GENERIC_PREFIX = "Couldn't clone the repository (git exited with code"


def classify_clone_error(exc: subprocess.CalledProcessError) -> str:
    """Map a failed `git clone`'s stderr to a curated, plain-English reason
    via CLONE_ERROR_KEYWORDS. Falls back to a generic message carrying only
    the exit code — never the raw stderr text itself."""
    stderr = (exc.stderr or "").lower()
    for keywords, message in CLONE_ERROR_KEYWORDS:
        if any(keyword in stderr for keyword in keywords):
            return message
    return f"{_CLONE_ERROR_GENERIC_PREFIX} {exc.returncode})."


def is_clone_error_message(message: str) -> bool:
    """True if `message` is one of classify_clone_error's curated outputs."""
    return (
        any(message == friendly for _, friendly in CLONE_ERROR_KEYWORDS)
        or message.startswith(_CLONE_ERROR_GENERIC_PREFIX)
    )


def _clone_git_repo(rag_source: RagSource, dest_dir: Path) -> None:
    """Shallow-clone `rag_source.git_url` into `dest_dir`.

    For a private repo, the PAT — `rag_source.auth_token` if set (an
    in-memory token resolved by the caller, e.g. kb_repo_service.
    descriptor_from_repo decrypting a self-service repo's stored token),
    otherwise read from `os.environ[rag_source.auth_token_env]` — is sent via
    the Azure-recommended HTTP Basic auth header — `-c
    http.extraHeader="Authorization: Basic <base64(':'+PAT)>"` — rather than
    embedded in the clone URL. The PAT itself is never logged: only the repo
    URL/ref are, and the constructed `cmd` list (which does contain the
    base64-encoded header value, as git requires) is passed straight to
    subprocess.run — it is never printed, logged, or otherwise persisted.
    """
    config_args: list[str] = []
    pat = rag_source.auth_token or (
        os.environ.get(rag_source.auth_token_env, "") if rag_source.auth_token_env else ""
    )
    if pat:
        token_b64 = base64.b64encode(f":{pat}".encode("utf-8")).decode("ascii")
        config_args = ["-c", f"http.extraHeader=Authorization: Basic {token_b64}"]

    # core.longpaths lets git check out paths longer than Windows' 260-char
    # MAX_PATH (a deeply-nested repo otherwise fails checkout with "Filename
    # too long"); it's a no-op/harmless on other platforms.
    cmd = ["git", "-c", "core.longpaths=true"] + config_args + ["clone", "--depth", "1"]
    if rag_source.git_ref:
        cmd += ["--branch", rag_source.git_ref]
    cmd += [rag_source.git_url, str(dest_dir)]

    log.info("[ingestion] git clone %s (ref=%s)", rag_source.git_url, rag_source.git_ref or "default")
    try:
        subprocess.run(cmd, check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as exc:
        # Deliberately do not include `cmd` (carries the auth header) or raw
        # stderr (git may echo back the failing URL) in the raised message —
        # only a curated, safe reason from classify_clone_error.
        raise RuntimeError(classify_clone_error(exc)) from None


# ── File walking ────────────────────────────────────────────────────────────


def _is_excluded(rel_posix: str, exclude_globs: list[str]) -> bool:
    return any(fnmatch.fnmatch(rel_posix, pat) for pat in exclude_globs)


def _iter_files(root_dirs: list[Path], include_globs: list[str], exclude_globs: list[str]):
    """Yield (root_dir, path, rel_label) for every file under root_dirs that
    matches include_globs and not exclude_globs. rel_label is a
    root-relative, forward-slash path used as the stable kb_documents.path
    key (so re-ingesting a freshly re-cloned repo — a new temp dir each time
    — still resolves to the same document row)."""
    seen: set[Path] = set()
    for root in root_dirs:
        if not root.exists():
            continue
        for pattern in include_globs:
            for path in root.glob(pattern):
                if not path.is_file():
                    continue
                resolved = path.resolve()
                if resolved in seen:
                    continue
                if any(part in _ALWAYS_SKIP_DIR_NAMES for part in path.relative_to(root).parts):
                    continue
                rel = path.relative_to(root).as_posix()
                if _is_excluded(rel, exclude_globs):
                    continue
                seen.add(resolved)
                yield root, path, f"{root.name}/{rel}"


# ── Text extraction ─────────────────────────────────────────────────────────


def _load_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        log.warning("[ingestion] pypdf not installed — skipping PDF %s", path.name)
        return ""
    try:
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        log.warning("[ingestion] failed reading PDF %s: %s", path.name, type(exc).__name__)
        return ""


def _load_text(path: Path) -> str:
    if path.suffix.lower() == _PDF_SUFFIX:
        return _load_pdf_text(path)
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as exc:
        log.warning("[ingestion] failed reading %s: %s", path.name, type(exc).__name__)
        return ""


def _sha256_of_file(path: Path) -> str:
    h = hashlib.sha256()
    h.update(path.read_bytes())
    return h.hexdigest()


# ── Upload text extraction (Phase 10) ───────────────────────────────────────
#
# extract_text(filename, data) is the in-memory counterpart to _load_text
# above — used by main.py's POST /kb-repos/{id}/files, which never writes an
# uploaded file to disk. Every parser here reads from an io.BytesIO wrapping
# the raw upload bytes, no temp files. SUPPORTED_UPLOAD_SUFFIXES is the
# allowlist the endpoint rejects against *before* calling extract_text —
# anything else (images, audio, video, archives, other binaries) is refused
# rather than silently producing empty/garbage text.

# Plain-text / markup / code — decoded as UTF-8 (errors="ignore" so a stray
# non-UTF-8 byte never fails the whole upload).
_UPLOAD_TEXT_SUFFIXES = {
    ".md", ".txt", ".rst", ".csv", ".json", ".yaml", ".yml", ".html", ".xml",
    ".log", ".py", ".ts", ".tsx", ".js", ".jsx", ".css", ".sql", ".ini",
    ".toml", ".cfg", ".sh", ".java", ".go", ".rb", ".php", ".c", ".cpp", ".h",
}
_UPLOAD_EXCEL_SUFFIXES = {".xlsx", ".xls"}
_UPLOAD_DOCX_SUFFIX = ".docx"
_UPLOAD_PPTX_SUFFIX = ".pptx"

SUPPORTED_UPLOAD_SUFFIXES: set[str] = (
    _UPLOAD_TEXT_SUFFIXES
    | _UPLOAD_EXCEL_SUFFIXES
    | {_PDF_SUFFIX, _UPLOAD_DOCX_SUFFIX, _UPLOAD_PPTX_SUFFIX}
)


def _extract_pdf_bytes(data: bytes) -> str:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        log.warning("[ingestion] failed reading uploaded PDF: %s", type(exc).__name__)
        return ""


def _extract_excel_bytes(data: bytes) -> str:
    import pandas as pd

    try:
        sheets = pd.read_excel(io.BytesIO(data), sheet_name=None)
    except Exception as exc:
        log.warning("[ingestion] failed reading uploaded spreadsheet: %s", type(exc).__name__)
        return ""
    parts = [f"# {name}\n{df.to_csv(index=False)}" for name, df in sheets.items()]
    return "\n\n".join(parts)


def _extract_docx_bytes(data: bytes) -> str:
    from docx import Document

    try:
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as exc:
        log.warning("[ingestion] failed reading uploaded docx: %s", type(exc).__name__)
        return ""


def _extract_pptx_bytes(data: bytes) -> str:
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
        slide_texts = []
        for slide in prs.slides:
            shape_texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            slide_texts.append("\n".join(shape_texts))
        return "\n\n".join(slide_texts)
    except Exception as exc:
        log.warning("[ingestion] failed reading uploaded pptx: %s", type(exc).__name__)
        return ""


def extract_text(filename: str, data: bytes) -> str:
    """Extract plain text from an in-memory uploaded file's bytes, dispatched
    by `filename`'s (lowercased) suffix. Raises ValueError for any suffix not
    in SUPPORTED_UPLOAD_SUFFIXES — callers (main.py's upload endpoint) are
    expected to check membership themselves first so they can report a
    friendly per-file rejection reason instead of relying on this exception,
    but it is raised regardless as a defensive backstop."""
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_UPLOAD_SUFFIXES:
        raise ValueError(f"Unsupported file type: {suffix or filename!r}")
    if suffix == _PDF_SUFFIX:
        return _extract_pdf_bytes(data)
    if suffix in _UPLOAD_EXCEL_SUFFIXES:
        return _extract_excel_bytes(data)
    if suffix == _UPLOAD_DOCX_SUFFIX:
        return _extract_docx_bytes(data)
    if suffix == _UPLOAD_PPTX_SUFFIX:
        return _extract_pptx_bytes(data)
    return data.decode("utf-8", errors="ignore")


# ── Chunking ─────────────────────────────────────────────────────────────────


def _chunk_text(text: str) -> list[str]:
    words = text.split()
    if not words:
        return []
    step = max(_CHUNK_WORDS - _CHUNK_OVERLAP_WORDS, 1)
    chunks: list[str] = []
    start = 0
    while start < len(words):
        window = words[start : start + _CHUNK_WORDS]
        chunks.append(" ".join(window))
        if start + _CHUNK_WORDS >= len(words):
            break
        start += step
    return chunks


# ── Ingestion ────────────────────────────────────────────────────────────────


def ingest_text_document(kb_id: str, path_label: str, text: str, *, sha256: str | None = None) -> int:
    """Chunk `text` (via `_chunk_text`), batch-embed it (`embeddings.
    embed_texts`), and upsert into the configured VectorStore as the
    `kb_documents` row at `path_label` within `kb_id` — creating that row if
    it's new, or replacing its prior chunks in place if it already exists
    (upsert_chunks only inserts; this deletes the stale chunks first). Shared
    by ingest_kb (one call per changed repo file, sha256 = the file's own
    content hash so a later re-ingest can skip it if unchanged) and main.py's
    upload endpoint (sha256 omitted — every upload always re-chunks/re-embeds
    in full; there is no byte-identical-skip concept for a one-shot upload).

    Returns the number of chunks written (0 for an empty document — the
    kb_documents row is still recorded)."""
    import db
    from models import KbChunk, KbDocument
    from sqlalchemy import select
    from vector_store import get_vector_store

    import embeddings

    sha = sha256 if sha256 is not None else hashlib.sha256(text.encode("utf-8")).hexdigest()
    chunk_texts = _chunk_text(text)

    with db.SyncSessionLocal() as session:
        existing_doc = session.execute(
            select(KbDocument).where(KbDocument.kb_id == kb_id, KbDocument.path == path_label)
        ).scalar_one_or_none()

        if existing_doc is None:
            existing_doc = KbDocument(kb_id=kb_id, path=path_label, sha256=sha)
            session.add(existing_doc)
            session.flush()  # assign existing_doc.id
        else:
            existing_doc.sha256 = sha
            existing_doc.updated_at = _utcnow()
            # Drop this document's stale chunks before re-embedding —
            # upsert_chunks only inserts, it doesn't replace.
            session.query(KbChunk).filter(KbChunk.document_id == existing_doc.id).delete()
        session.commit()

        if not chunk_texts:
            return 0  # e.g. an empty file — document recorded, no chunks

        vectors = embeddings.embed_texts(chunk_texts)
        chunk_dicts = [
            {
                "document_id": existing_doc.id,
                "chunk_index": i,
                "text": t,
                "source_ref": f"{path_label}#{i}",
                "embedding": vec,
            }
            for i, (t, vec) in enumerate(zip(chunk_texts, vectors))
        ]
        get_vector_store().upsert_chunks(kb_id, chunk_dicts)
        return len(chunk_dicts)


_UPLOADS_PATH_PREFIX = "uploads/"


def ingest_kb(descriptor: KBDescriptor) -> dict:
    """Ingest `descriptor`'s rag source into the vector store. See module
    docstring for the full flow. Returns {"documents", "chunks", "skipped",
    "removed"}."""
    rag_source = _resolve_rag_source(descriptor)
    counts = {"documents": 0, "chunks": 0, "skipped": 0, "removed": 0}
    if rag_source is None:
        log.debug("[ingestion] kb=%s has no rag source configured — nothing to ingest", descriptor.id)
        return counts

    import db
    from models import KbChunk, KbDocument
    from sqlalchemy import select

    tmp_dir: Path | None = None
    try:
        if rag_source.git_url:
            tmp_dir = Path(tempfile.mkdtemp(prefix="rule_agent_rag_"))
            _clone_git_repo(rag_source, tmp_dir)
            root_dirs = [tmp_dir / r for r in rag_source.roots] if rag_source.roots else [tmp_dir]
        else:
            root_dirs = [_resolve_root(r) for r in rag_source.roots]

        include_globs = rag_source.include_globs or ["**/*"]
        exclude_globs = rag_source.exclude_globs or []

        matched_paths: set[str] = set()
        with db.SyncSessionLocal() as session:
            for _root, path, rel_label in _iter_files(root_dirs, include_globs, exclude_globs):
                matched_paths.add(rel_label)
                sha = _sha256_of_file(path)

                existing_doc = session.execute(
                    select(KbDocument).where(
                        KbDocument.kb_id == descriptor.id, KbDocument.path == rel_label
                    )
                ).scalar_one_or_none()

                if existing_doc is not None and existing_doc.sha256 == sha:
                    counts["skipped"] += 1
                    continue

                text = _load_text(path)
                counts["documents"] += 1
                counts["chunks"] += ingest_text_document(descriptor.id, rel_label, text, sha256=sha)

            # Prune vanished files — see module docstring point 5. Only a
            # successful walk reaches here (a failed clone raises earlier).
            stale_docs = session.execute(
                select(KbDocument).where(KbDocument.kb_id == descriptor.id)
            ).scalars().all()
            for doc in stale_docs:
                if doc.path in matched_paths or doc.path.startswith(_UPLOADS_PATH_PREFIX):
                    continue
                session.query(KbChunk).filter(KbChunk.document_id == doc.id).delete()
                session.delete(doc)
                counts["removed"] += 1
            session.commit()

        return counts
    finally:
        if tmp_dir is not None:
            shutil.rmtree(tmp_dir, ignore_errors=True)
