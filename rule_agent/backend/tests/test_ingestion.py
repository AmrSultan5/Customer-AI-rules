"""
Tests for ingestion.py — ingest_kb() over local files (chunk + embed +
upsert, sha256 idempotency across re-ingest) and the Azure Repos git-clone
auth path. The embedder is mocked (embeddings._embed_batch_fn) and git is
mocked (ingestion.subprocess.run) throughout — no real OpenAI/network/git
calls are made anywhere in this file.
"""

import base64
import subprocess
from pathlib import Path

import pytest

import embeddings
import ingestion
from kb._schema import KBDescriptor
from vector_store import get_vector_store


def _fake_embed_batch(texts: list[str]) -> list[list[float]]:
    # Deterministic, cheap "embeddings" derived from text length so results
    # are stable within a test run without any real embedding model.
    return [[float(len(t) % 7), float(len(t) % 5), 1.0] for t in texts]


@pytest.fixture(autouse=True)
def _mock_embedder(monkeypatch):
    monkeypatch.setattr(embeddings, "_embed_batch_fn", _fake_embed_batch)


def _rag_descriptor(kb_id: str, roots: list[str], **rag_kwargs) -> KBDescriptor:
    return KBDescriptor(
        id=kb_id,
        name="Test RAG KB",
        description="A rag-only test KB for ingestion tests.",
        adapter="rag",
        retrieval_mode="rag",
        source={"kind": "rag", "roots": roots, **rag_kwargs},
        id_pattern=r"\b([A-Z]{2,8}_\d+)\b",
        vocab={"entity_singular": "doc", "entity_plural": "docs"},
        prompts={"repository_label": "Docs"},
    )


def _structured_only_descriptor(kb_id: str) -> KBDescriptor:
    return KBDescriptor(
        id=kb_id,
        name="Structured only",
        description="No rag source at all.",
        adapter="structured",
        retrieval_mode="structured",
        source={"kind": "structured", "files": {}, "dirs": {}},
        id_pattern=r"\b([A-Z]{2,8}_\d+)\b",
        vocab={"entity_singular": "x", "entity_plural": "xs"},
        prompts={"repository_label": "X"},
    )


# ── Local-file ingestion: chunk + embed + upsert ────────────────────────────


def test_ingest_kb_chunks_embeds_and_upserts_local_files(tmp_path):
    (tmp_path / "a.md").write_text("Hello world. " * 50, encoding="utf-8")
    (tmp_path / "b.txt").write_text("Another document about testing.", encoding="utf-8")
    descriptor = _rag_descriptor("test_ingest_local", [str(tmp_path)])

    counts = ingestion.ingest_kb(descriptor)

    assert counts["documents"] == 2
    assert counts["chunks"] >= 2
    assert counts["skipped"] == 0

    store = get_vector_store()
    assert store.count(descriptor.id) == counts["chunks"]


def test_ingest_kb_no_rag_source_is_a_noop():
    descriptor = _structured_only_descriptor("test_ingest_structured_only")
    counts = ingestion.ingest_kb(descriptor)
    assert counts == {"documents": 0, "chunks": 0, "skipped": 0, "removed": 0}


# ── sha256 idempotency across re-ingest ─────────────────────────────────────


def test_reingest_skips_unchanged_files(tmp_path):
    (tmp_path / "a.md").write_text("Stable content that does not change.", encoding="utf-8")
    descriptor = _rag_descriptor("test_ingest_idempotent", [str(tmp_path)])

    first = ingestion.ingest_kb(descriptor)
    assert first["documents"] == 1
    assert first["skipped"] == 0
    assert first["chunks"] >= 1

    second = ingestion.ingest_kb(descriptor)
    assert second["documents"] == 0
    assert second["skipped"] == 1
    assert second["chunks"] == 0

    store = get_vector_store()
    assert store.count(descriptor.id) == first["chunks"]  # nothing duplicated


def _kb_document_paths(kb_id: str) -> set[str]:
    import db
    from models import KbDocument
    from sqlalchemy import select

    with db.SyncSessionLocal() as session:
        return set(session.execute(
            select(KbDocument.path).where(KbDocument.kb_id == kb_id)
        ).scalars().all())


def test_reingest_prunes_deleted_file_and_its_chunks(tmp_path):
    a = tmp_path / "a.md"
    b = tmp_path / "b.md"
    a.write_text("First document, kept around.", encoding="utf-8")
    b.write_text("Second document, about to be deleted, with enough distinct words.", encoding="utf-8")
    descriptor = _rag_descriptor("test_ingest_prune_deleted", [str(tmp_path)])

    first = ingestion.ingest_kb(descriptor)
    assert first["documents"] == 2
    assert first["removed"] == 0

    store = get_vector_store()
    assert store.count(descriptor.id) == first["chunks"]

    b.unlink()
    second = ingestion.ingest_kb(descriptor)
    assert second["removed"] == 1
    assert second["skipped"] == 1  # a.md unchanged, so it's skipped not re-embedded

    paths = _kb_document_paths(descriptor.id)
    assert paths == {f"{tmp_path.name}/a.md"}

    # b.md's chunks are gone too — nothing left over from the vanished doc.
    assert store.count(descriptor.id) < first["chunks"]
    assert store.count(descriptor.id) > 0  # a.md's own chunks remain


def test_reingest_preserves_uploads_documents(tmp_path):
    (tmp_path / "a.md").write_text("Repo document content.", encoding="utf-8")
    kb_id = "test_ingest_prune_uploads"
    descriptor = _rag_descriptor(kb_id, [str(tmp_path)])

    ingestion.ingest_kb(descriptor)
    ingestion.ingest_text_document(kb_id, "uploads/manual.md", "Manually uploaded content.")

    store = get_vector_store()
    before = store.count(kb_id)
    assert before > 0

    # Re-ingest the SAME repo file set — the uploaded doc must survive since
    # it's never part of the repo walk.
    result = ingestion.ingest_kb(descriptor)
    assert result["removed"] == 0

    assert "uploads/manual.md" in _kb_document_paths(kb_id)


def test_reingest_narrower_glob_prunes_now_excluded_files(tmp_path):
    (tmp_path / "a.md").write_text("Markdown document.", encoding="utf-8")
    (tmp_path / "b.txt").write_text("Text document that will be excluded next.", encoding="utf-8")
    kb_id = "test_ingest_prune_glob_narrow"
    wide = _rag_descriptor(kb_id, [str(tmp_path)], include_globs=["**/*"])

    first = ingestion.ingest_kb(wide)
    assert first["documents"] == 2

    narrow = _rag_descriptor(kb_id, [str(tmp_path)], include_globs=["**/*.md"])
    second = ingestion.ingest_kb(narrow)
    assert second["removed"] == 1

    assert _kb_document_paths(kb_id) == {f"{tmp_path.name}/a.md"}


def test_reingest_reembeds_changed_file_and_replaces_its_chunks(tmp_path):
    f = tmp_path / "a.md"
    f.write_text("Version one of the document.", encoding="utf-8")
    descriptor = _rag_descriptor("test_ingest_changed", [str(tmp_path)])

    first = ingestion.ingest_kb(descriptor)
    assert first["documents"] == 1

    f.write_text("Version two, with substantially different content that changes the sha256 hash.", encoding="utf-8")
    second = ingestion.ingest_kb(descriptor)

    assert second["documents"] == 1
    assert second["skipped"] == 0

    store = get_vector_store()
    # Old chunks were replaced, not accumulated alongside the new ones.
    assert store.count(descriptor.id) == second["chunks"]


# ── Azure DevOps Repos git-clone auth path (subprocess fully mocked) ───────


def test_git_clone_uses_http_extra_header_basic_auth_and_never_logs_the_pat(monkeypatch, tmp_path, caplog):
    captured: dict = {}

    def fake_run(cmd, check, capture_output, text):
        captured["cmd"] = cmd
        dest = Path(cmd[-1])
        dest.mkdir(parents=True, exist_ok=True)
        (dest / "doc.md").write_text("Cloned content from the mocked repo.", encoding="utf-8")

        class _Result:
            returncode = 0
            stdout = ""
            stderr = ""

        return _Result()

    monkeypatch.setattr(ingestion.subprocess, "run", fake_run)
    monkeypatch.setenv("TEST_REPO_PAT", "super-secret-pat-value")

    descriptor = _rag_descriptor(
        "test_ingest_git",
        roots=[],
        git_url="https://dev.azure.com/org/project/_git/repo",
        git_ref="main",
        auth_token_env="TEST_REPO_PAT",
    )

    with caplog.at_level("DEBUG"):
        counts = ingestion.ingest_kb(descriptor)

    assert counts["documents"] == 1
    assert counts["chunks"] >= 1

    cmd = captured["cmd"]

    # The bare repo URL is passed as-is — no embedded credentials.
    url_arg = next(c for c in cmd if c == descriptor.source.git_url)
    assert "super-secret-pat-value" not in url_arg
    assert "@dev.azure.com" not in url_arg

    # The PAT is sent via -c http.extraHeader="Authorization: Basic <b64>",
    # never the URL — this is the point of the mechanism, so the header
    # argument legitimately carries the (base64-encoded) token.
    header_arg = next(c for c in cmd if c.startswith("http.extraHeader="))
    token_b64 = base64.b64encode(b":super-secret-pat-value").decode("ascii")
    assert header_arg == f"http.extraHeader=Authorization: Basic {token_b64}"

    # The raw PAT must never appear in any emitted log record.
    for record in caplog.records:
        assert "super-secret-pat-value" not in record.getMessage()


def test_git_clone_without_pat_omits_auth_header(monkeypatch, tmp_path):
    captured: dict = {}

    def fake_run(cmd, check, capture_output, text):
        captured["cmd"] = cmd
        dest = Path(cmd[-1])
        dest.mkdir(parents=True, exist_ok=True)
        (dest / "doc.md").write_text("Public repo content.", encoding="utf-8")

        class _Result:
            returncode = 0

        return _Result()

    monkeypatch.setattr(ingestion.subprocess, "run", fake_run)
    monkeypatch.delenv("PUBLIC_REPO_PAT", raising=False)

    descriptor = _rag_descriptor(
        "test_ingest_git_public",
        roots=[],
        git_url="https://dev.azure.com/org/project/_git/public-repo",
        auth_token_env="PUBLIC_REPO_PAT",
    )

    ingestion.ingest_kb(descriptor)

    cmd = captured["cmd"]
    assert not any(c.startswith("http.extraHeader=") for c in cmd)
    # No auth header when there's no PAT, but core.longpaths is always set.
    assert "core.longpaths=true" in cmd


# ── classify_clone_error / is_clone_error_message (friendly classification) ─


@pytest.mark.parametrize(
    "stderr, expected",
    [
        (
            "remote: Authentication failed for 'https://dev.azure.com/org/project/_git/repo'",
            "Couldn't access the repository — check the access token (private repos require one).",
        ),
        (
            "fatal: could not read Username for 'https://dev.azure.com': terminal prompts disabled",
            "Couldn't access the repository — check the access token (private repos require one).",
        ),
        (
            "remote: TF401019: The Git repository with name or identifier does not exist",
            "Repository not found — check the Git URL and branch.",
        ),
        (
            "remote: Repository not found.",
            "Repository not found — check the Git URL and branch.",
        ),
        (
            "fatal: unable to access 'https://dev.azure.com/org/project/_git/repo/': Could not resolve host: dev.azure.com",
            "Couldn't reach the repository host — check the URL and your network.",
        ),
        (
            "fatal: unable to access 'https://dev.azure.com/x': Connection timed out",
            "Couldn't reach the repository host — check the URL and your network.",
        ),
        (
            # git prefixes path errors with "unable to access '<path>':" too, so
            # this must NOT be misread as a network failure.
            "warning: unable to access 'a/b/c/.gitattributes': Filename too long\n"
            "fatal: cannot create directory at 'a/b/c/processes': Filename too long",
            "The repository downloaded, but its files couldn't be unpacked here — usually a path that's too long for this system (a deeply-nested folder). Try a repo with shorter paths.",
        ),
        (
            "fatal: unable to access 'https://x/': SSL certificate problem: unable to get local issuer certificate",
            "Couldn't verify the repository's security certificate — usually a corporate network/proxy. Ask IT to trust the proxy's certificate for git.",
        ),
    ],
)
def test_classify_clone_error_maps_known_keywords_to_friendly_reason(stderr, expected):
    exc = subprocess.CalledProcessError(returncode=128, cmd=["git", "clone"], output="", stderr=stderr)
    assert ingestion.classify_clone_error(exc) == expected


def test_classify_clone_error_falls_back_to_generic_message_with_exit_code():
    exc = subprocess.CalledProcessError(returncode=17, cmd=["git", "clone"], output="", stderr="some other git failure")
    message = ingestion.classify_clone_error(exc)
    assert message == "Couldn't clone the repository (git exited with code 17)."
    # The raw stderr text must never leak into the curated message.
    assert "some other git failure" not in message


def test_is_clone_error_message_recognizes_curated_and_generic_messages():
    assert ingestion.is_clone_error_message(
        "Couldn't access the repository — check the access token (private repos require one)."
    )
    assert ingestion.is_clone_error_message("Couldn't clone the repository (git exited with code 42).")
    assert not ingestion.is_clone_error_message("Some unrelated RuntimeError text.")


def test_clone_git_repo_raises_friendly_message_never_raw_stderr_or_pat(monkeypatch, tmp_path):
    def fake_run(cmd, check, capture_output, text):
        raise subprocess.CalledProcessError(
            returncode=128,
            cmd=cmd,
            output="",
            stderr="remote: Authentication failed for 'https://dev.azure.com/org/project/_git/repo'",
        )

    monkeypatch.setattr(ingestion.subprocess, "run", fake_run)
    monkeypatch.setenv("CLONE_FAIL_PAT", "super-secret-pat-value")

    from kb._schema import RagSource

    rag_source = RagSource(
        roots=[],
        git_url="https://dev.azure.com/org/project/_git/repo",
        auth_token_env="CLONE_FAIL_PAT",
    )

    with pytest.raises(RuntimeError) as exc_info:
        ingestion._clone_git_repo(rag_source, tmp_path / "dest")

    message = str(exc_info.value)
    assert message == "Couldn't access the repository — check the access token (private repos require one)."
    assert "super-secret-pat-value" not in message
    assert "Authentication failed" not in message


# ── extract_text / SUPPORTED_UPLOAD_SUFFIXES (Phase 10 uploads) ────────────


def test_supported_upload_suffixes_covers_expected_types():
    expected = {
        ".md", ".txt", ".rst", ".csv", ".json", ".yaml", ".yml", ".html", ".xml",
        ".log", ".py", ".ts", ".tsx", ".js", ".jsx", ".css", ".sql", ".ini",
        ".toml", ".cfg", ".sh", ".java", ".go", ".rb", ".php", ".c", ".cpp", ".h",
        ".pdf", ".xlsx", ".xls", ".docx", ".pptx",
    }
    assert ingestion.SUPPORTED_UPLOAD_SUFFIXES == expected


def test_extract_text_plain_text_suffix_decodes_utf8():
    text = ingestion.extract_text("notes.md", "Hello **world**".encode("utf-8"))
    assert text == "Hello **world**"


def test_extract_text_ignores_undecodable_bytes():
    # A stray non-UTF-8 byte must not blow up extraction (errors="ignore").
    data = b"before" + bytes([0xFF]) + b"after"
    text = ingestion.extract_text("notes.txt", data)
    assert "before" in text and "after" in text


def test_extract_text_xlsx_concatenates_sheets():
    import io

    import pandas as pd

    buf = io.BytesIO()
    pd.DataFrame({"a": [1, 2], "b": ["x", "y"]}).to_excel(buf, index=False, sheet_name="Sheet1")
    text = ingestion.extract_text("data.xlsx", buf.getvalue())
    assert "a" in text and "b" in text
    assert "x" in text and "y" in text


def test_extract_text_docx():
    import io

    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello from a Word document.")
    buf = io.BytesIO()
    doc.save(buf)
    text = ingestion.extract_text("report.docx", buf.getvalue())
    assert "Hello from a Word document." in text


def test_extract_text_pptx_per_slide():
    import io

    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Slide Title Text"
    buf = io.BytesIO()
    prs.save(buf)
    text = ingestion.extract_text("deck.pptx", buf.getvalue())
    assert "Slide Title Text" in text


def test_extract_text_pdf_via_pypdf():
    import io

    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    # A blank page has no extractable text — assert extraction runs cleanly
    # (never raises) and returns a string, rather than asserting content.
    text = ingestion.extract_text("blank.pdf", buf.getvalue())
    assert isinstance(text, str)


def test_extract_text_unsupported_suffix_is_rejected():
    assert ".png" not in ingestion.SUPPORTED_UPLOAD_SUFFIXES
    with pytest.raises(ValueError):
        ingestion.extract_text("photo.png", b"\x89PNG\r\n")


# ── ingest_text_document (shared chunk/embed/upsert helper) ─────────────────


def test_ingest_text_document_chunks_embeds_and_upserts():
    kb_id = "test_upload_helper"
    n = ingestion.ingest_text_document(kb_id, "uploads/a.md", "Some uploaded content here.")
    assert n >= 1

    store = get_vector_store()
    assert store.count(kb_id) == n


def test_ingest_text_document_replaces_prior_chunks_at_same_path():
    kb_id = "test_upload_helper_replace"
    n1 = ingestion.ingest_text_document(kb_id, "uploads/a.md", "Version one content.")
    store = get_vector_store()
    assert store.count(kb_id) == n1

    n2 = ingestion.ingest_text_document(kb_id, "uploads/a.md", "Version two, replaced content that differs.")
    # Old chunks were replaced, not accumulated alongside the new ones.
    assert store.count(kb_id) == n2


def test_ingest_text_document_empty_text_records_document_with_zero_chunks():
    kb_id = "test_upload_helper_empty"
    n = ingestion.ingest_text_document(kb_id, "uploads/empty.md", "")
    assert n == 0
    assert get_vector_store().count(kb_id) == 0
